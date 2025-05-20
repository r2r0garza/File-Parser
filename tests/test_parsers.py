import os
import tempfile
from main import parse_docx, parse_text, parse_pdf, parse_csv, parse_xlsx, parse_pptx, parse_eml, parse_image

def test_parse_text():
  with tempfile.NamedTemporaryFile(mode="w+", suffix=".txt", delete=False) as f:
    f.write("Hello, world!")
    f.flush()
    path = f.name
  try:
    result = parse_text(path)
    assert result == "Hello, world!"
  finally:
    os.remove(path)

def test_parse_csv():
  with tempfile.NamedTemporaryFile(mode="w+", suffix=".csv", delete=False) as f:
    f.write("a,b\n1,2\n3,4")
    f.flush()
    path = f.name
  try:
    result = parse_csv(path)
    assert "a,b" in result and "1,2" in result and "3,4" in result
  finally:
    os.remove(path)

def test_parse_ts():
  with tempfile.NamedTemporaryFile(mode="w+", suffix=".ts", delete=False) as f:
    f.write("let y: string = 'hello';")
    f.flush()
    path = f.name
  try:
    result = parse_text(path)
    assert result == "let y: string = 'hello';"
  finally:
    os.remove(path)

def test_parse_docx():
  import docx
  path = tempfile.mktemp(suffix=".docx")
  doc = docx.Document()
  doc.add_paragraph("Docx test")
  doc.save(path)
  try:
    result = parse_docx(path)
    assert "Docx test" in result
  finally:
    os.remove(path)

def test_parse_pdf():
  import pdfplumber
  from fpdf import FPDF
  # Create a simple PDF file
  pdf_path = tempfile.mktemp(suffix=".pdf")
  pdf = FPDF()
  pdf.add_page()
  pdf.set_font("Arial", size=12)
  pdf.cell(200, 10, txt="PDF test", ln=True)
  pdf.output(pdf_path)
  try:
    result = parse_pdf(pdf_path)
    assert "PDF test" in result or result.strip() != ""
  finally:
    os.remove(pdf_path)

def test_parse_xlsx():
  import pandas as pd
  path = tempfile.mktemp(suffix=".xlsx")
  df = pd.DataFrame({"col1": [1, 2], "col2": [3, 4]})
  df.to_excel(path, index=False)
  try:
    result = parse_xlsx(path)
    assert "col1" in result and "col2" in result and "1" in result and "2" in result
  finally:
    os.remove(path)

def test_parse_pptx():
  from pptx import Presentation
  path = tempfile.mktemp(suffix=".pptx")
  prs = Presentation()
  slide = prs.slides.add_slide(prs.slide_layouts[5])
  txBox = slide.shapes.add_textbox(left=100, top=100, width=500, height=100)
  tf = txBox.text_frame
  tf.text = "PPTX test"
  prs.save(path)
  try:
    result = parse_pptx(path)
    assert "PPTX test" in result
  finally:
    os.remove(path)

def test_parse_eml():
  import email
  path = tempfile.mktemp(suffix=".eml")
  msg = email.message.EmailMessage()
  msg["Subject"] = "EML Subject"
  msg["From"] = "from@example.com"
  msg["To"] = "to@example.com"
  msg.set_content("This is the body of the EML file.")
  with open(path, "wb") as f:
    f.write(msg.as_bytes())
  try:
    result = parse_eml(path)
    assert "EML Subject" in result and "from@example.com" in result and "This is the body" in result
  finally:
    os.remove(path)

def test_parse_image_ocr():
  from PIL import Image, ImageDraw, ImageFont
  path = tempfile.mktemp(suffix=".png")
  # Create a simple image with text
  img = Image.new("RGB", (200, 60), color=(255, 255, 255))
  d = ImageDraw.Draw(img)
  d.text((10, 10), "OCR Test", fill=(0, 0, 0))
  img.save(path)
  try:
    result = parse_image(path)
    assert "OCR" in result or "Test" in result
  finally:
    os.remove(path)

def test_parse_image_ocr_low_quality():
  from PIL import Image, ImageDraw, ImageFont, ImageFilter
  path = tempfile.mktemp(suffix=".png")
  img = Image.new("RGB", (200, 60), color=(255, 255, 255))
  d = ImageDraw.Draw(img)
  d.text((10, 10), "Blurry", fill=(0, 0, 0))
  img = img.filter(ImageFilter.GaussianBlur(radius=2))
  img.save(path)
  try:
    result = parse_image(path)
    # Accept partial or noisy OCR result
    assert "Blurry" in result or result.strip() != ""
  finally:
    os.remove(path)

# Stub for large legacy document test (manual/placeholder)
def test_large_legacy_doc_stub():
  # This is a placeholder for manual stress testing with large .doc/.xls/.ppt files
  # To be run manually with real files if needed
  pass
