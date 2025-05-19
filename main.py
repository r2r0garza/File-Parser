import logging
import base64
import requests
from fastapi import FastAPI, UploadFile, File, HTTPException, Request, Form
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import os
import tempfile
import docx
import pdfplumber
import pandas as pd
import openpyxl
import pptx
import extract_msg
import email
import email.policy
import subprocess
import pytesseract
from PIL import Image
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# LLaVA configuration
LLAVA_USE = os.getenv("LLAVA_USE", "false").lower() == "true"
LLAVA_URL = os.getenv("LLAVA_URL", "http://localhost:1234/v1/chat/completions")
LLAVA_MODEL_NAME = os.getenv("LLAVA_MODEL_NAME", "llava-1.6-mistral-7b")

app = FastAPI()

# Enable CORS for all origins (for development; restrict in production)
app.add_middleware(
  CORSMiddleware,
  allow_origins=["*"],
  allow_credentials=True,
  allow_methods=["*"],
  allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
  level=logging.INFO,
  format="%(asctime)s [%(levelname)s] %(message)s"
)
logger = logging.getLogger(__name__)

@app.get("/")
def read_root():
  return {"message": "Document Parser API is running."}

@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
  logger.error(f"Unhandled error: {exc}")
  return JSONResponse(
    status_code=500,
    content={"detail": "Internal server error."}
  )

def get_file_size(file_path: str) -> int:
  try:
    return os.path.getsize(file_path)
  except Exception:
    return 0

def get_pdf_page_count(file_path: str) -> int:
  try:
    with pdfplumber.open(file_path) as pdf:
      return len(pdf.pages)
  except Exception:
    return 0

def get_pptx_slide_count(file_path: str) -> int:
  try:
    prs = pptx.Presentation(file_path)
    return len(prs.slides)
  except Exception:
    return 0

def get_csv_shape(file_path: str) -> dict:
  try:
    df = pd.read_csv(file_path)
    return {"rows": df.shape[0], "columns": df.shape[1]}
  except Exception:
    return {"rows": 0, "columns": 0}

def get_xlsx_shape(file_path: str) -> dict:
  try:
    df = pd.read_excel(file_path, engine="openpyxl")
    return {"rows": df.shape[0], "columns": df.shape[1]}
  except Exception:
    return {"rows": 0, "columns": 0}

# File type detection utility
def detect_file_type(filename: str) -> str:
  ext = os.path.splitext(filename)[1].lower()
  if ext.startswith("."):
    ext = ext[1:]
  return ext

# Parser for .docx files
def parse_docx(file_path: str) -> str:
  try:
    doc = docx.Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text
  except Exception as e:
    logger.error(f"Error parsing .docx: {e}")
    return ""

# Parser for legacy .doc, .xls, .ppt files using unoconv + libreoffice
def parse_legacy_office(file_path: str, target_ext: str) -> str:
  try:
    out_path = file_path + ".converted" + target_ext
    subprocess.run(["unoconv", "-f", target_ext.lstrip("."), "-o", out_path, file_path], check=True)
    if target_ext == ".docx":
      return parse_docx(out_path)
    elif target_ext == ".xlsx":
      return parse_xlsx(out_path)
    elif target_ext == ".pptx":
      return parse_pptx(out_path)
    else:
      return ""
  except Exception as e:
    logger.error(f"Error parsing legacy office file: {e}")
    return "Legacy format parsing requires unoconv/libreoffice installed."

# Parser for text files (.txt, .md, .log)
def parse_text(file_path: str) -> str:
  try:
    with open(file_path, "r", encoding="utf-8") as f:
      return f.read()
  except Exception as e:
    logger.error(f"Error parsing text file: {e}")
    return ""

# Parser for .pdf files with OCR fallback
def parse_pdf(file_path: str) -> str:
  try:
    text = ""
    with pdfplumber.open(file_path) as pdf:
      for page in pdf.pages:
        page_text = page.extract_text()
        if page_text and page_text.strip():
          text += page_text
        else:
          # OCR fallback for image-based PDF pages
          img = page.to_image(resolution=300).original
          ocr_text = pytesseract.image_to_string(img)
          text += ocr_text
        text += "\n"
    return text.strip()
  except Exception as e:
    logger.error(f"Error parsing .pdf: {e}")
    return ""

# Parser for .csv files
def parse_csv(file_path: str) -> str:
  try:
    df = pd.read_csv(file_path)
    return df.to_csv(index=False)
  except Exception as e:
    logger.error(f"Error parsing .csv: {e}")
    return ""

# Parser for .xlsx files
def parse_xlsx(file_path: str) -> str:
  try:
    df = pd.read_excel(file_path, engine="openpyxl")
    return df.to_csv(index=False)
  except Exception as e:
    logger.error(f"Error parsing .xlsx: {e}")
    return ""

# Parser for .pptx files
def parse_pptx(file_path: str) -> str:
  try:
    prs = pptx.Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
      for shape in slide.shapes:
        if hasattr(shape, "text"):
          text_runs.append(shape.text)
    return "\n".join(text_runs)
  except Exception as e:
    logger.error(f"Error parsing .pptx: {e}")
    return ""

# Parser for .eml files
def parse_eml(file_path: str) -> str:
  try:
    with open(file_path, "rb") as f:
      msg = email.message_from_binary_file(f, policy=email.policy.default)
    subject = msg.get("subject", "")
    from_ = msg.get("from", "")
    to = msg.get("to", "")
    body = ""
    if msg.is_multipart():
      for part in msg.walk():
        if part.get_content_type() == "text/plain":
          body += part.get_content()
    else:
      body = msg.get_content()
    return f"Subject: {subject}\nFrom: {from_}\nTo: {to}\n\n{body}"
  except Exception as e:
    logger.error(f"Error parsing .eml: {e}")
    return ""

# Parser for .msg files (Outlook)
def parse_msg(file_path: str) -> str:
  try:
    msg = extract_msg.Message(file_path)
    subject = msg.subject or ""
    sender = msg.sender or ""
    to = msg.to or ""
    body = msg.body or ""
    return f"Subject: {subject}\nFrom: {sender}\nTo: {to}\n\n{body}"
  except Exception as e:
    logger.error(f"Error parsing .msg: {e}")
    return ""

# Parser for image files (.png, .jpg, .jpeg) using pytesseract or LLaVA
def parse_image(file_path: str) -> str:
  try:
    # If LLaVA is not enabled, use pytesseract
    if not LLAVA_USE:
      img = Image.open(file_path)
      text = pytesseract.image_to_string(img)
      return text
    else:
      # For LLaVA, we'll return a placeholder as the actual processing
      # happens in the /caption endpoint
      return "[Image processed by LLaVA. Use /caption endpoint for detailed description.]"
  except Exception as e:
    logger.error(f"Error parsing image file: {e}")
    return ""

# Router for dispatching to parsers
def parse_file_router(file_path: str, filetype: str) -> str:
  if filetype == "docx":
    return parse_docx(file_path)
  elif filetype == "doc":
    return parse_legacy_office(file_path, ".docx")
  elif filetype == "xls":
    return parse_legacy_office(file_path, ".xlsx")
  elif filetype == "ppt":
    return parse_legacy_office(file_path, ".pptx")
  elif filetype in {"txt", "md", "log"}:
    return parse_text(file_path)
  elif filetype == "pdf":
    return parse_pdf(file_path)
  elif filetype == "csv":
    return parse_csv(file_path)
  elif filetype == "xlsx":
    return parse_xlsx(file_path)
  elif filetype == "pptx":
    return parse_pptx(file_path)
  elif filetype == "eml":
    return parse_eml(file_path)
  elif filetype == "msg":
    return parse_msg(file_path)
  elif filetype in {"png", "jpg", "jpeg"}:
    return parse_image(file_path)
  else:
    logger.warning(f"Unsupported file type: {filetype}")
    return ""

def extract_metadata(file_path: str, filetype: str) -> dict:
  meta = {"size_bytes": get_file_size(file_path)}
  if filetype == "pdf":
    meta["page_count"] = get_pdf_page_count(file_path)
  elif filetype == "pptx":
    meta["slide_count"] = get_pptx_slide_count(file_path)
  elif filetype == "csv":
    meta.update(get_csv_shape(file_path))
  elif filetype == "xlsx":
    meta.update(get_xlsx_shape(file_path))
  return meta

# /parse endpoint for file uploads
@app.post("/parse")
async def parse_upload(file: UploadFile = File(...)):
  try:
    suffix = os.path.splitext(file.filename)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
      content = await file.read()
      tmp.write(content)
      tmp_path = tmp.name

    filetype = detect_file_type(file.filename)
    parsed_content = parse_file_router(tmp_path, filetype)
    metadata = extract_metadata(tmp_path, filetype)

    os.remove(tmp_path)

    return {
      "filename": file.filename,
      "filetype": filetype,
      "metadata": metadata,
      "content": parsed_content
    }
  except Exception as e:
    logger.error(f"Error in /parse: {e}")
    raise HTTPException(status_code=500, detail=f"Failed to parse file: {str(e)}")

# Pydantic model for /parse-path
class ParsePathRequest(BaseModel):
  filepath: str

@app.post("/parse-path")
async def parse_path(req: ParsePathRequest):
  try:
    if not os.path.isfile(req.filepath):
      raise HTTPException(status_code=404, detail="File not found.")

    filetype = detect_file_type(req.filepath)
    parsed_content = parse_file_router(req.filepath, filetype)
    metadata = extract_metadata(req.filepath, filetype)

    return {
      "filename": os.path.basename(req.filepath),
      "filetype": filetype,
      "metadata": metadata,
      "content": parsed_content
    }
  except Exception as e:
    logger.error(f"Error in /parse-path: {e}")
    raise HTTPException(status_code=500, detail=f"Failed to parse file path: {str(e)}")

# Utility: DataFrame to Markdown
def dataframe_to_markdown(df):
  if df.empty:
    return ""
  header = "| " + " | ".join(map(str, df.columns)) + " |"
  separator = "| " + " | ".join(["---"] * len(df.columns)) + " |"
  rows = ["| " + " | ".join(map(str, row)) + " |" for row in df.values]
  return "\n".join([header, separator] + rows)

# /xlsx-to-md endpoint
from fastapi.responses import PlainTextResponse

@app.post("/xlsx-to-md", response_class=PlainTextResponse)
async def xlsx_to_markdown(file: UploadFile = File(...)):
  try:
    suffix = os.path.splitext(file.filename)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
      content = await file.read()
      tmp.write(content)
      tmp_path = tmp.name

    # Read all sheets
    xls = pd.ExcelFile(tmp_path, engine="openpyxl")
    md_parts = []
    for sheet_name in xls.sheet_names:
      df = xls.parse(sheet_name)
      md = dataframe_to_markdown(df)
      if md.strip():
        md_parts.append(f"## {sheet_name}\n\n{md}\n")
    os.remove(tmp_path)

    if not md_parts:
      return "No data found in the XLSX file."
    return "\n".join(md_parts)
  except Exception as e:
    logger.error(f"Error in /xlsx-to-md: {e}")
    raise HTTPException(status_code=500, detail=f"Failed to convert XLSX to Markdown: {str(e)}")

# /caption endpoint for LLaVA image processing
@app.post("/caption")
async def caption(
  image: UploadFile = File(...),
  system_prompt: str = Form("You are a helpful assistant.")
):
  try:
    if not LLAVA_USE:
      raise HTTPException(
        status_code=400, 
        detail="LLaVA integration is disabled. Set LLAVA_USE=true in .env to enable."
      )
    
    # Read and encode image to base64
    image_content = await image.read()
    image_b64 = base64.b64encode(image_content).decode("utf-8")
    image_mime = image.content_type or "image/png"
    
    # Build payload for LLaVA
    payload = {
      "model": LLAVA_MODEL_NAME,
      "messages": [
        {"role": "system", "content": system_prompt},
        {
          "role": "user",
          "content": [
            {"type": "text", "text": "Describe the image in detail."},
            {
              "type": "image_url",
              "image_url": {
                "url": f"data:{image_mime};base64,{image_b64}"
              }
            }
          ]
        }
      ],
      "max_tokens": 512
    }
    
    # Send request to LLaVA server
    response = requests.post(LLAVA_URL, json=payload)
    
    # Check if the request was successful
    if response.status_code != 200:
      logger.error(f"LLaVA server error: {response.text}")
      raise HTTPException(
        status_code=response.status_code,
        detail=f"LLaVA server error: {response.text}"
      )
    
    return response.json()
  except Exception as e:
    logger.error(f"Error in /caption: {e}")
    raise HTTPException(status_code=500, detail=f"Failed to process image with LLaVA: {str(e)}")
